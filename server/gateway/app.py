"""team-brain gateway — single public front door for the brain.

Routes (all bearer-gated except /healthz):
  /healthz              liveness (no auth)
  /ingest/upload        POST multipart: extract -> chunk -> store original -> push to agentmemory
  /docs/{doc_id}        GET original file back
  /docs                 GET manifest (list of ingested docs)
  /viewer, /viewer/*    reverse-proxy to the agentmemory viewer (:3113); browser
                        login via /viewer?key=<secret> (sets an HttpOnly cookie)
  /agentmemory/*        reverse-proxy to the agentmemory REST API (:3111); accepts
                        bearer OR the viewer cookie (gateway injects the engine bearer)

Design notes:
  - Extraction is server-side so client VMs stay thin (just upload bytes).
  - Each chunk is stored in agentmemory as type=reference with provenance in the
    text (doc id, source filename, section) so smart-search returns it alongside
    everything else, and /docs/{id} serves the original for deep reads.
  - sha256 of file bytes = idempotency: re-uploading an unchanged file is a no-op.
"""
import asyncio
import hashlib
import io
import json
import os
import re
import sqlite3
import time
from pathlib import Path

import httpx
from fastapi import FastAPI, Header, HTTPException, Request, UploadFile, File, Form
from fastapi.responses import (JSONResponse, Response, FileResponse,
                               StreamingResponse, RedirectResponse)

AM_BASE = os.environ.get("AM_BASE", "http://localhost:3111")
VIEWER_BASE = os.environ.get("VIEWER_BASE", "http://localhost:3113")
DOCS_DIR = Path(os.environ.get("BRAIN_DOCS_DIR", str(Path.home() / "brain-docs")))
DB_PATH = DOCS_DIR / "manifest.db"
CHUNK_TARGET = int(os.environ.get("BRAIN_CHUNK_CHARS", "1500"))


def _secret() -> str:
    s = os.environ.get("AGENTMEMORY_SECRET", "").strip()
    if s:
        return s
    f = Path.home() / ".agentmemory" / "team_secret.txt"
    return f.read_text().strip() if f.exists() else ""


SECRET = _secret()
DOCS_DIR.mkdir(parents=True, exist_ok=True)
# Disable built-in Swagger/OpenAPI: it would shadow our /docs route and is an
# unauthenticated info-leak surface on a public URL.
app = FastAPI(title="team-brain gateway", docs_url=None, redoc_url=None, openapi_url=None)


def db():
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        """CREATE TABLE IF NOT EXISTS docs(
             id TEXT PRIMARY KEY, sha256 TEXT UNIQUE, filename TEXT, ext TEXT,
             bytes INTEGER, chunks INTEGER, user TEXT, note TEXT, created REAL)"""
    )
    return conn


def require_auth(authorization: str | None):
    if not SECRET:
        return  # unconfigured: fail open only if no secret set (dev)
    if authorization != f"Bearer {SECRET}":
        raise HTTPException(status_code=401, detail="unauthorized")


# Browser access to the viewer can't send a bearer header, so the viewer + the
# /agentmemory/* calls its JS makes also accept an HttpOnly cookie set via
# /viewer?key=<secret>. The cookie carries the shared secret (admin/ops-grade — fine
# for who reaches the dashboard; per-user gating would be a later thing).
VIEWER_COOKIE = "brain_viewer"


def authed(request: Request, authorization: str | None) -> bool:
    if not SECRET:
        return True
    return (authorization == f"Bearer {SECRET}"
            or request.cookies.get(VIEWER_COOKIE) == SECRET)


def require_auth_browser(request: Request, authorization: str | None):
    if not authed(request, authorization):
        raise HTTPException(status_code=401, detail="unauthorized")


# ---------- extraction ----------
def extract_text(data: bytes, ext: str, filename: str) -> str:
    ext = ext.lower().lstrip(".")
    if ext in ("md", "markdown", "txt", "text", "rst", "csv", "json", "yaml", "yml"):
        return data.decode("utf-8", errors="replace")
    if ext == "pdf":
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        parts = []
        for i, page in enumerate(doc):
            t = page.get_text("text").strip()
            if t:
                parts.append(f"[page {i+1}]\n{t}")
        return "\n\n".join(parts)
    if ext == "docx":
        import docx
        d = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in d.paragraphs if p.text.strip())
    if ext in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = [sh.text for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
            if texts:
                parts.append(f"[slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    # fallback: best-effort decode
    return data.decode("utf-8", errors="replace")


def chunk_text(text: str, target: int = CHUNK_TARGET) -> list[str]:
    """Structure-aware-ish: split on blank lines / [page|slide] markers, pack to ~target."""
    blocks = re.split(r"\n\s*\n", text)
    chunks, cur = [], ""
    for b in blocks:
        b = b.strip()
        if not b:
            continue
        if len(cur) + len(b) + 2 > target and cur:
            chunks.append(cur)
            cur = b
        else:
            cur = f"{cur}\n\n{b}" if cur else b
    if cur:
        chunks.append(cur)
    return chunks


async def am_post(path: str, payload: dict, timeout: float = 30) -> httpx.Response:
    async with httpx.AsyncClient(timeout=timeout) as c:
        return await c.post(f"{AM_BASE}{path}", json=payload,
                            headers={"Authorization": f"Bearer {SECRET}"})


# ---------- cataloger / maintenance ----------
# Spin-down-native scheduling: instead of a clock cron (the box is suspended at
# 3am), we run the cataloger when the brain is already awake and enough new
# material has accumulated. A manual/external POST /maintenance/run is also
# exposed so a Podclave-side scheduler can force a run on demand.
MAINT = {"writes": 0, "last_run": 0.0, "running": False, "last_result": None}
MAINT_WRITES = int(os.environ.get("BRAIN_MAINT_WRITES", "20"))      # trigger after N writes
MAINT_MIN_SECS = int(os.environ.get("BRAIN_MAINT_MIN_SECS", "1800"))  # ...but at most this often
_maint_lock = asyncio.Lock()


SPRITE_SOCK = os.environ.get("SPRITE_SOCK", "/.sprite/api.sock")
KEEPALIVE = "brain-consolidating"


async def _sprite_api(method: str, path: str, payload: dict | None = None):
    """Call the local Sprite tasks API over its unix socket (keep-alive)."""
    transport = httpx.AsyncHTTPTransport(uds=SPRITE_SOCK)
    async with httpx.AsyncClient(transport=transport, base_url="http://sprite", timeout=10) as c:
        return await c.request(method, path, json=payload)


async def run_maintenance(reason: str = "manual") -> dict:
    """Run the cataloger: consolidate -> reflect -> auto-forget. One at a time.
    Holds a Sprite keep-alive task so the box can't auto-suspend mid-catalog."""
    if _maint_lock.locked():
        return {"skipped": "already_running"}
    async with _maint_lock:
        MAINT["running"] = True
        out = {"reason": reason, "started": time.time()}
        try:
            try:
                await _sprite_api("POST", "/v1/tasks", {"name": KEEPALIVE, "expire": "30m"})
            except Exception:  # noqa: BLE001
                pass  # keep-alive is best-effort; never block the catalog on it
            for step in ("consolidate-pipeline", "reflect", "auto-forget"):
                try:
                    r = await am_post(f"/agentmemory/{step}", {}, timeout=170)
                    out[step] = r.json() if r.headers.get("content-type", "").startswith("application/json") else r.status_code
                except Exception as e:  # noqa: BLE001
                    out[step] = f"error: {e}"
            out["finished"] = time.time()
        finally:
            try:
                await _sprite_api("DELETE", f"/v1/tasks/{KEEPALIVE}")
            except Exception:  # noqa: BLE001
                pass
            MAINT.update(running=False, last_run=time.time(), writes=0, last_result=out)
        return out


def note_writes(n: int = 1):
    """Count writes and kick off maintenance in the background when due."""
    MAINT["writes"] += n
    due = (MAINT["writes"] >= MAINT_WRITES
           and (time.time() - MAINT["last_run"]) >= MAINT_MIN_SECS
           and not MAINT["running"])
    if due:
        asyncio.create_task(run_maintenance(reason="activity"))


# ---------- endpoints ----------
@app.get("/healthz")
async def healthz():
    return {"status": "ok", "service": "team-brain-gateway"}


@app.post("/ingest/upload")
async def ingest_upload(
    file: UploadFile = File(...),
    note: str = Form(""),
    user: str = Form("unknown"),
    authorization: str | None = Header(default=None),
):
    require_auth(authorization)
    data = await file.read()
    sha = hashlib.sha256(data).hexdigest()
    ext = Path(file.filename or "").suffix.lstrip(".") or "txt"

    conn = db()
    row = conn.execute("SELECT id, chunks FROM docs WHERE sha256=?", (sha,)).fetchone()
    if row:
        conn.close()
        return {"status": "already_ingested", "doc_id": row[0], "chunks": row[1],
                "filename": file.filename}

    doc_id = sha[:16]
    text = extract_text(data, ext, file.filename or "")
    if not text.strip():
        raise HTTPException(status_code=422, detail="no extractable text")
    chunks = chunk_text(text)

    # store original for /docs deep reads
    (DOCS_DIR / f"{doc_id}.{ext}").write_bytes(data)

    pushed = 0
    for idx, ch in enumerate(chunks):
        # Provenance goes LAST so the chunk's real content dominates the embedding
        # (a leading metadata prefix measurably hurt semantic recall).
        prov = f"[source: {file.filename} | doc:{doc_id} | chunk {idx+1}/{len(chunks)} | filed by {user}"
        prov += f" | note: {note}]" if note else "]"
        body = f"{ch}\n\n{prov}"
        r = await am_post("/agentmemory/remember",
                          {"content": body, "type": "reference"})
        if r.status_code in (200, 201):
            pushed += 1

    conn.execute(
        "INSERT INTO docs VALUES (?,?,?,?,?,?,?,?,?)",
        (doc_id, sha, file.filename, ext, len(data), len(chunks), user, note, time.time()),
    )
    conn.commit()
    conn.close()
    note_writes(pushed)
    return {"status": "ingested", "doc_id": doc_id, "filename": file.filename,
            "ext": ext, "chunks": len(chunks), "chunks_stored": pushed,
            "deep_read": f"/docs/{doc_id}"}


@app.get("/docs")
async def list_docs(authorization: str | None = Header(default=None)):
    require_auth(authorization)
    conn = db()
    rows = conn.execute(
        "SELECT id, filename, ext, chunks, user, note, created FROM docs ORDER BY created DESC"
    ).fetchall()
    conn.close()
    return {"docs": [
        {"doc_id": r[0], "filename": r[1], "ext": r[2], "chunks": r[3],
         "user": r[4], "note": r[5], "created": r[6]} for r in rows]}


@app.get("/docs/{doc_id}")
async def get_doc(doc_id: str, authorization: str | None = Header(default=None)):
    require_auth(authorization)
    conn = db()
    row = conn.execute("SELECT filename, ext FROM docs WHERE id=?", (doc_id,)).fetchone()
    conn.close()
    if not row:
        raise HTTPException(status_code=404, detail="unknown doc_id")
    path = DOCS_DIR / f"{doc_id}.{row[1]}"
    if not path.exists():
        raise HTTPException(status_code=410, detail="original not stored")
    return FileResponse(path, filename=row[0])


# ---------- reverse proxies ----------
async def _proxy(base: str, path: str, request: Request, inject_auth: bool = False) -> Response:
    url = f"{base}/{path}"
    body = await request.body()
    headers = {k: v for k, v in request.headers.items()
               if k.lower() not in ("host", "content-length")}
    if inject_auth and SECRET:
        # Gateway authenticates to the engine itself, regardless of how the client
        # authed to the gateway (bearer OR viewer cookie) — so browser/cookie requests
        # still satisfy the engine's own bearer check upstream.
        headers["authorization"] = f"Bearer {SECRET}"
    async with httpx.AsyncClient(timeout=60) as c:
        r = await c.request(request.method, url, content=body, headers=headers,
                            params=request.query_params)
    resp_headers = {k: v for k, v in r.headers.items()
                    if k.lower() not in ("content-encoding", "transfer-encoding", "content-length", "connection")}
    return Response(content=r.content, status_code=r.status_code, headers=resp_headers)


@app.post("/maintenance/run")
async def maintenance_run(authorization: str | None = Header(default=None)):
    require_auth(authorization)
    return await run_maintenance(reason="manual")


@app.get("/maintenance/status")
async def maintenance_status(authorization: str | None = Header(default=None)):
    require_auth(authorization)
    return {"writes_since_run": MAINT["writes"], "running": MAINT["running"],
            "last_run": MAINT["last_run"], "trigger_at_writes": MAINT_WRITES,
            "min_interval_secs": MAINT_MIN_SECS, "last_result": MAINT["last_result"]}


# ---------- remember passthrough (write-count for the cataloger) ----------
# We deliberately do NOT dedup here. agentmemory supersedes near-duplicates NATIVELY on
# write (Jaccard >=0.7 -> marks the old memory isLatest=false, bumps version, keeps the
# *newer* phrasing). Our previous token-set dedup pre-empted that path and kept the
# OLDER text (and made supersedes look broken), so it's removed. This handler only
# forwards and counts the write so the activity-triggered cataloger still fires.
# (Interactive saves now go via the agentmemory MCP -> /agentmemory/mcp/call, which the
# catch-all proxies; those aren't counted here, but the scheduled /maintenance/run
# covers consolidation regardless.)
@app.post("/agentmemory/remember")
async def remember(request: Request, authorization: str | None = Header(default=None)):
    require_auth(authorization)
    body = await request.json()
    r = await am_post("/agentmemory/remember", body)
    note_writes(1)
    try:
        return JSONResponse(r.json(), status_code=r.status_code)
    except Exception:  # noqa: BLE001
        return Response(content=r.content, status_code=r.status_code)


@app.api_route("/agentmemory/{path:path}",
               methods=["GET", "POST", "PUT", "DELETE", "PATCH"])
async def proxy_am(path: str, request: Request,
                   authorization: str | None = Header(default=None)):
    require_auth_browser(request, authorization)  # bearer OR viewer cookie
    # inject_auth so the viewer's cookie-authed XHRs still satisfy the engine bearer.
    return await _proxy(f"{AM_BASE}/agentmemory", path, request, inject_auth=True)


@app.get("/favicon.svg")
async def favicon(request: Request):
    return await _proxy(VIEWER_BASE, "favicon.svg", request)  # cosmetic; unauth


@app.api_route("/viewer", methods=["GET"])
@app.api_route("/viewer/{path:path}", methods=["GET"])
async def proxy_viewer(request: Request, path: str = "",
                       authorization: str | None = Header(default=None)):
    # Browser login: /viewer?key=<secret> sets an HttpOnly cookie and redirects
    # (303) to a clean URL so the secret never lingers in the address bar / history.
    key = request.query_params.get("key")
    if key is not None:
        if SECRET and key != SECRET:
            raise HTTPException(status_code=401, detail="bad key")
        dest = "/viewer/" + path if path else "/viewer"
        r = RedirectResponse(url=dest, status_code=303)
        r.set_cookie(VIEWER_COOKIE, SECRET, httponly=True, secure=True,
                     samesite="lax", max_age=86400, path="/")
        return r
    require_auth_browser(request, authorization)
    return await _proxy(VIEWER_BASE, path, request, inject_auth=True)
